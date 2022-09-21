import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from urllib.request import urlopen
from bs4 import BeautifulSoup

def TransVol(vol):
    if (vol == "창세기")or(vol == "창"):
        return"gen"
    elif (vol == "출애굽기")or(vol == "출"):
        return"exo"
    elif (vol == "레위기")or(vol == "레"):
        return"lev"
    elif (vol == "민수기")or(vol == "민"):
        return"num"
    elif (vol == "신명기")or(vol == "신"):
        return"deu"
    elif (vol == "여호수아")or(vol == "수"):
        return"jos"
    elif (vol == "사사기")or(vol == "삿"):
        return"jdg"
    elif (vol == "룻")or(vol == "룻"):
        return"rut"
    elif (vol == "사무엘상")or(vol == "삼상"):
        return"1sa"
    elif (vol == "사무엘하")or(vol == "삼하"):
        return"2sa"
    elif (vol == "열왕기상")or(vol == "왕상"):
        return"1ki"
    elif (vol == "열왕기하")or(vol == "왕하"):
        return"2ki"
    elif (vol == "역대상")or(vol == "대상"):
        return"1ch"
    elif (vol == "역대하")or(vol == "대하"):
        return"2ch"
    elif (vol == "에스라")or(vol == "스"):
        return"ezr"
    elif (vol == "느헤미야")or(vol == "느"):
        return"neh"
    elif (vol == "에스더")or(vol == "에"):
        return"est"
    elif (vol == "욥")or(vol == "욥"):
        return"job"
    elif (vol == "시편")or(vol == "시"):
        return"psa"
    elif (vol == "잠언")or(vol == "잠"):
        return"pro"
    elif (vol == "전도서")or(vol == "전"):
        return"ecc"
    elif (vol == "아가")or(vol == "아"):
        return"sng"
    elif (vol == "이사야")or(vol == "사"):
        return"isa"
    elif (vol == "예레미야")or(vol == "렘"):
        return"jer"
    elif (vol == "예레미야 애가")or(vol == "애"):
        return"lam"
    elif (vol == "에스겔")or(vol == "겔"):
        return"ezk"
    elif (vol == "다니엘")or(vol == "단"):
        return"dan"
    elif (vol == "호세아")or(vol == "호"):
        return"hos"
    elif (vol == "요엘")or(vol == "욜"):
        return"jol"
    elif (vol == "아모스")or(vol == "암"):
        return"amo"
    elif (vol == "오바댜")or(vol == "옵"):
        return"oba"
    elif (vol == "요나")or(vol == "욘"):
        return"jnh"
    elif (vol == "미가")or(vol == "미"):
        return"mic"
    elif (vol == "나훔")or(vol == "나"):
        return"nam"
    elif (vol == "하박국")or(vol == "합"):
        return"hab"
    elif (vol == "스바냐")or(vol == "습"):
        return"zep"
    elif (vol == "학개")or(vol == "학"):
        return"hag"
    elif (vol == "스가랴")or(vol == "슥"):
        return"zec"
    elif (vol == "말라기")or(vol == "말"):
        return"mal"
    elif (vol == "마태복음")or(vol == "마"):
        return"mat"
    elif (vol == "마가복음")or(vol == "막"):
        return"mrk"
    elif (vol == "누가복음")or(vol == "눅"):
        return"luk"
    elif (vol == "요한복음")or(vol == "요"):
        return"jhn"
    elif (vol == "사도행전")or(vol == "행"):
        return"act"
    elif (vol == "로마서")or(vol == "롬"):
        return"rom"
    elif (vol == "고린도전서")or(vol == "고전"):
        return"1co"
    elif (vol == "고린도후서")or(vol == "고후"):
        return"2co"
    elif (vol == "갈라디아")or(vol == "갈"):
        return"gal"
    elif (vol == "에베소서")or(vol == "엡"):
        return"eph"
    elif (vol == "빌립보서")or(vol == "빌"):
        return"php"
    elif (vol == "골로새서")or(vol == "골"):
        return"col"
    elif (vol == "데살로니가전서")or(vol == "살전"):
        return"1th"
    elif (vol == "데살로니가후서")or(vol == "살후"):
        return"2th"
    elif (vol == "디모데전서")or(vol == "딤전"):
        return"1ti"
    elif (vol == "디모데후서")or(vol == "딤후"):
        return"2ti"
    elif (vol == "디도서")or(vol == "딛"):
        return"tit"
    elif (vol == "빌레몬서")or(vol == "몬"):
        return"phm"
    elif (vol == "히브리서")or(vol == "히"):
        return"heb"
    elif (vol == "야고보서")or(vol == "약"):
        return"jas"
    elif (vol == "베드로전서")or(vol == "벧전"):
        return"1pe"
    elif (vol == "베드로후서")or(vol == "벧후"):
        return"2pe"
    elif (vol == "요한1서")or(vol == "요일"):
        return"1jn"
    elif (vol == "요한2서")or(vol == "요이"):
        return"2jn"
    elif (vol == "요한3서")or(vol == "요삼"):
        return"3jn"
    elif (vol == "유다서")or(vol == "유 "):
        return"jud"
    elif (vol == "요한계시록")or(vol == "계 "):
        return"rev"


def GenBible(prs_name, save_name):
    prs = Presentation(prs_name)

    vol = input("성경 권 이름을 입력하세요: ")
    chap = int(input("성경 장 수를 입력하세요: "))
    start = int(input("시작 절을 입력하세요: "))
    end = int(input("끝 절을 입력하세요: "))
    vol = TransVol(vol)
    url = "http://bible.godpia.com/read/reading.asp?ver=gae&ver2=&vol=" + vol + "&chap=" + str(chap) + "&sec="
    html = urlopen(url)  
    bsObject = BeautifulSoup(html, "html.parser")

    for i in range(start,end+1):
        sec = '#gae_'+vol+'_'+str(chap)+'_'+str(i)+' > span'
        line = bsObject.select_one(sec)
        line = line.get_text()
        line = line.replace(str(i), str(i)+'. ')
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = line

    prs.save(save_name)

def GenBibleAdditional(prs_name, save_name, file_name):
    prs = Presentation(prs_name)
    f = open(file_name, 'r')
    lines = f.readlines()
    first = 1
    for line in lines:
        if first:
            l_title = line
            first = 0
        else:
            title_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = l_title
            subtitle.text = line

    f.close()
    prs.save(save_name)