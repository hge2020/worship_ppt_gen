import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR


def GenAny(prs_name, save_name, filename):
    prs = Presentation(prs_name)

    f = open(filename, 'r')
    lines = f.readlines()
    for line in lines:
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = line

    f.close()

    prs.save(save_name)