import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR


def GenLyrics(prs_name, save_name, file_name):
    prs = Presentation(prs_name)

    f = open(file_name, 'r')
    lines = f.readlines()
    second = 0
    l_title = 0

    for line in lines:
        line = line.strip()
        if line == "!!":
            l_title = 1
        elif l_title:
            title_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = line
            title.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.TEXT_2
            l_title = 0
            second = 0
        elif not second:
            l = line
            second = 1
        else:
            l = l + '\n' + line
            title_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = l
            second = 0

    f.close()
    prs.save(save_name)